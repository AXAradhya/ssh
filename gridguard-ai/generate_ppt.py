import os
import subprocess
import sys

def install(package):
    print(f"Installing {package}...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    install('python-pptx')
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

prs = Presentation()

# Function to add a dark background to match GridGuard's UI
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 14, 31) # #0a0e1f

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    set_dark_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 212, 255) # Cyan
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    return slide

def add_content_slide(prs, title_text, content_bullets):
    slide_layout = prs.slide_layouts[1] # Content Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    set_dark_background(slide)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 212, 255) # Cyan
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for idx, bullet in enumerate(content_bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.font.size = Pt(20)

slides_data = [
    {
        "type": "title",
        "title": "GridGuard AI | The Next-Generation Smart Grid Platform",
        "subtitle": "Securing, Optimizing, and Decarbonizing the Future of Energy"
    },
    {
        "type": "content",
        "title": "Why The World Needs GridGuard AI",
        "content": [
            "Climate Volatility: Traditional grids cannot handle unpredictable load swings from extreme weather and massive influx of renewables (Solar/Wind).",
            "Cybersecurity Breaches: Rapid digitization exposes critical infrastructure to devastating DDoS attacks and rogue IP intrusions.",
            "Aging Assets: Multi-million dollar transformers are operating past rated lifespans, gambling with catastrophic 'loss-of-life' failures."
        ]
    },
    {
        "type": "content",
        "title": "The Omni-Aware Energy OS",
        "content": [
            "GridGuard AI is a live, self-healing Energy Operating System powered by 10 advanced Machine Learning engines.",
            "AI-Driven Resilience: Instantly identifies failures and re-routes power using self-healing topography.",
            "Military-Grade Security: Built-in ML Intrusion Detection System (IDS) blacklists malicious traffic in real-time.",
            "LLM 'Omni-Aware' Assistant: The only platform featuring a generative AI chatbot querying live sensors for instantaneous structural analysis."
        ]
    },
    {
        "type": "content",
        "title": "Real-Time GIS & Topology Mapping",
        "content": [
            "Substation Visibility with Pinpoint Accuracy via Interactive Live Leaflet Mapping.",
            "Deep-Tech Tracking: Moves beyond static dashboards using precision GPS coordinates (Lat/Lng) to map critical assets like the Bawana Solar Plant.",
            "Node-Level Telemetry: Clicking physical assets reveals live flow rates, structural capacity limits, and instant load/generation readings.",
            "Aesthetic Dominance: Rendered in a stunning, custom dark-mode cyber aesthetic."
        ]
    },
    {
        "type": "content",
        "title": "Predictive Maintenance via IEEE Physics",
        "content": [
            "Transformer Digital Twin: Stop failures before they black out a city.",
            "The Technology: Uses a structural physics engine to simulate internal 'Top-Oil' and 'Hot-Spot' core temperatures.",
            "Loss-of-Life Tracking: Automatically tracks the IEEE C57.91 Insulation Loss-of-Life percentage.",
            "Automated Action: Automatically schedules autonomous drone dispatch flights for physical inspection if health drops."
        ]
    },
    {
        "type": "content",
        "title": "Wholesale Bidding & VPP Aggregation",
        "content": [
            "Multi-Agent Economic Optimization: Turning the grid into a profit center.",
            "Virtual Power Plants (VPP): Aggregates thousands of residential roof-top solar panels and Tesla powerwalls into a single, massive virtual battery.",
            "Wholesale Bidding Agent: AI autonomously decides whether to STORE energy during low-cost hours, or SELL renewables during peak demand spikes."
        ]
    },
    {
        "type": "content",
        "title": "Talk to Your Grid. The Omni-Aware Engine",
        "content": [
            "Google Gemini LLM Integration for real-time conversational intelligence.",
            "The Capability: Operators ask 'What is the structural health of Rohini Substation and should we fly drones tomorrow?'",
            "The Magic: Instantly polls the VPP, IDS Security Scanner, PMU Disturbance monitor, and Digital Twin.",
            "Aggregates a live 'Omni-Context' payload to generate flawless, data-backed answers in seconds."
        ]
    },
    {
        "type": "content",
        "title": "Secure, Persistent, Scalable Architecture",
        "content": [
            "Strict SQL Persistence: Robust SQLite/SQLAlchemy embedded database permanently audits every scheduled drone and mitigated cyber attack.",
            "Zero-Trust Authentication: Wrapped entirely in heavily vaulted JSON Web Tokens (JWT) with strict role-based access control.",
            "React + FastAPI: Built utilizing the bleeding edge of modern web development (React 18 + Python FastAPI).",
            "Ensures sub-millisecond API response times even under heavy load."
        ]
    },
    {
        "type": "content",
        "title": "Scalability & The Future: What Comes Next?",
        "content": [
            "Decentralized Carbon Ledger: Securely logs utilized green energy into an immutable SQL ledger to automate carbon-credit generation.",
            "EV Fleet Orchestration: Actively tracking electric vehicles, paying drivers to discharge power into the grid during blackouts.",
            "Global Zone Expansion: Pluggable architecture allows instant additions of new cities, nations, or interconnected power pools."
        ]
    },
    {
        "type": "title",
        "title": "The Future is Connected.",
        "subtitle": "GridGuard AI isn't just a prototype—it is a fully engineered, deployable command center ready to protect tomorrow's infrastructure today.\n\n- Live Demonstration\n- Pilot Program Request\n- Q&A Session"
    }
]

for slide_data in slides_data:
    if slide_data["type"] == "title":
        add_title_slide(prs, slide_data["title"], slide_data["subtitle"])
    else:
        add_content_slide(prs, slide_data["title"], slide_data["content"])

output_path = os.path.join(os.getcwd(), 'GridGuard_AI_Presentation.pptx')
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
